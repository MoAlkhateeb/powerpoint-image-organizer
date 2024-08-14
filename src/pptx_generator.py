from pathlib import Path
from typing import Optional, List

from pptx.slide import Slide
from pptx import Presentation

from image_layout_manager import ImageLayoutManager
from pptx_settings import PPTXSettings, Pathlike, BLANK_SLIDE_LAYOUT


class PPTXGenerator:
    def __init__(self, settings: Optional[PPTXSettings] = None):
        self.presentation = None
        self.settings = settings or PPTXSettings()
        self.slide_size = None

    def create_presentation(
        self, presentation_path: Pathlike, override: bool = False
    ) -> None:
        if Path(presentation_path).exists() and not override:
            self.presentation = Presentation(str(presentation_path))
        else:
            self.presentation = Presentation()

        self.slide_size = (
            self.presentation.slide_width,
            self.presentation.slide_height,
        )

    def add_images(self, images: List[Pathlike]) -> None:
        images_grouped_by_four = [images[i : i + 4] for i in range(0, len(images), 4)]

        for group in images_grouped_by_four:
            slide = self._create_empty_slide()
            self._arrange_images(slide, group)

    def save_presentation(self, presentation_path: Pathlike) -> None:
        self.presentation.save(str(presentation_path))

    def _create_empty_slide(self) -> Slide:
        return self.presentation.slides.add_slide(
            self.presentation.slide_layouts[BLANK_SLIDE_LAYOUT]
        )

    def _arrange_images(self, slide: Slide, images: List[Pathlike]) -> None:
        if not images:
            return

        layout_manager = ImageLayoutManager(self.slide_size, self.settings)

        match len(images):
            case 1:
                layout_manager.add_single_image(slide, images[0])
            case 2:
                layout_manager.add_two_images(slide, images)
            case 3:
                layout_manager.add_three_images(slide, images)
            case 4:
                layout_manager.add_four_images(slide, images)
            case _:
                raise ValueError("Only up to 4 images can be added to a slide")


if __name__ == "__main__":
    pptx_generator = PPTXGenerator()
    pptx_generator.create_presentation("presentation.pptx", override=False)

    pptx_generator.settings.rounded = True
    pptx_generator.settings.color = "springgreen"
    images = list(Path("images").glob("group_1/*.jpg"))

    pptx_generator.add_images(images)

    pptx_generator.save_presentation("presentation2.pptx")
