from pptx.slide import Slide
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from typing import List, Tuple, Optional
from pptx_settings import PPTXSettings, Pathlike


class ImageLayoutManager:
    def __init__(self, slide_size, settings: Optional[PPTXSettings] = None):
        self.slide_width, self.slide_height = slide_size
        self.settings = settings or PPTXSettings()

    def _calculate_dimensions(
        self, num_columns: int, num_rows: int
    ) -> Tuple[float, float]:
        width = (
            self.slide_width
            - self.settings.left_margin
            - self.settings.right_margin
            - (num_columns - 1) * self.settings.h_center_margin
        ) / num_columns

        height = (
            self.slide_height
            - self.settings.top_margin
            - self.settings.bottom_margin
            - (num_rows - 1) * self.settings.v_center_margin
        ) / num_rows
        return width, height

    def _get_position(
        self, row: int, col: int, width: float, height: float
    ) -> Tuple[float, float]:
        x = self.settings.left_margin + col * (width + self.settings.h_center_margin)
        y = self.settings.top_margin + row * (height + self.settings.v_center_margin)
        return x, y

    def _add_image(
        self,
        slide: Slide,
        image_path: Pathlike,
        left: Inches,
        top: Inches,
        width: Inches,
        height: Inches,
    ) -> None:
        pic = slide.shapes.add_picture(str(image_path), left, top, width, height)
        if self.settings.rounded:
            pic.auto_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
        if self.settings.line_width > 0:
            pic.line.color.rgb = RGBColor(*self.settings.color)
            pic.line.width = self.settings.line_width

    def add_single_image(self, slide: Slide, image_path: Pathlike) -> None:
        width, height = self._calculate_dimensions(1, 1)
        self._add_image(
            slide,
            image_path,
            self.settings.left_margin,
            self.settings.top_margin,
            width,
            height,
        )

    def add_two_images(self, slide: Slide, image_paths: List[Pathlike]) -> None:
        width, height = self._calculate_dimensions(2, 1)
        for i, image_path in enumerate(image_paths):
            x, y = self._get_position(0, i, width, height)
            self._add_image(slide, image_path, x, y, width, height)

    def add_three_images(self, slide: Slide, image_paths: List[Pathlike]) -> None:
        width, height = self._calculate_dimensions(2, 2)
        for i, image_path in enumerate(image_paths):
            if i < 2:
                x, y = self._get_position(0, i, width, height)
            else:
                x = (
                    self.settings.left_margin
                    + width / 2
                    + self.settings.h_center_margin / 2
                )
                y = self.settings.top_margin + height + self.settings.v_center_margin
            self._add_image(slide, image_path, x, y, width, height)

    def add_four_images(self, slide: Slide, image_paths: List[Pathlike]) -> None:
        width, height = self._calculate_dimensions(2, 2)
        for i, image_path in enumerate(image_paths):
            row, col = divmod(i, 2)
            x, y = self._get_position(row, col, width, height)
            self._add_image(slide, image_path, x, y, width, height)
